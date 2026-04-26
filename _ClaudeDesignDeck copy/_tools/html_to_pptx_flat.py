"""
Flat-image PPTX export of site-deck v6.html.

Each slide is rendered to a 1920x1080 PNG via headless Chromium, then
dropped full-bleed onto a 16:9 PPTX slide. Pixel-faithful to the original;
not editable below the slide level (each slide is one image).

Run:
    python3 html_to_pptx_flat.py
"""
from __future__ import annotations

import http.server
import io
import os
import socketserver
import threading
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

DECK_DIR = Path(__file__).resolve().parent.parent
DECK_FILE = "site-deck v6.html"
OUTPUT = DECK_DIR / "exports" / "site-deck-v6-flat.pptx"

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080


class _QuietHandler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, format, *args):  # noqa
        pass


def start_server(directory: Path, port: int = 8765) -> socketserver.TCPServer:
    os.chdir(directory)
    socketserver.TCPServer.allow_reuse_address = True
    httpd = socketserver.TCPServer(("127.0.0.1", port), _QuietHandler)
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return httpd


def main() -> None:
    httpd = start_server(DECK_DIR)
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            ctx = browser.new_context(
                viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
                device_scale_factor=2,  # crisper PNGs for big monitors / projection
            )
            page = ctx.new_page()
            url = f"http://127.0.0.1:8765/{DECK_FILE.replace(' ', '%20')}"
            page.goto(url, wait_until="networkidle")

            # Wait for fonts and the deck-stage component to be live.
            page.wait_for_function("document.fonts && document.fonts.ready")
            page.wait_for_function(
                "document.querySelector('deck-stage') "
                "&& document.querySelector('deck-stage')._slides "
                "&& document.querySelector('deck-stage')._slides.length > 0"
            )

            # Hide the deck-stage's prev/next overlay buttons so they don't
            # appear in the rendered screenshots.
            page.add_style_tag(content="""
              deck-stage::part(nav), deck-stage .nav, deck-stage .overlay,
              deck-stage button.prev, deck-stage button.next, deck-stage .btn {
                display: none !important; visibility: hidden !important;
              }
            """)

            total = page.evaluate("document.querySelector('deck-stage')._slides.length")
            print(f"deck-stage reports {total} slides")

            prs = Presentation()
            prs.slide_width = Emu(12192000)   # 13.333"  → 1920px
            prs.slide_height = Emu(6858000)   #  7.5"    → 1080px
            blank = prs.slide_layouts[6]

            for i in range(total):
                page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
                # Let the slide settle: arrival animation, image loads, font swap.
                page.wait_for_timeout(400)
                page.wait_for_load_state("networkidle")
                page.wait_for_function("document.fonts.ready")

                stage = page.query_selector("deck-stage")
                png = stage.screenshot(type="png", omit_background=False)
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(
                    io.BytesIO(png), Emu(0), Emu(0),
                    width=prs.slide_width, height=prs.slide_height,
                )
                print(f"slide {i:02d}: rendered {len(png) // 1024} KB")

            OUTPUT.parent.mkdir(parents=True, exist_ok=True)
            prs.save(OUTPUT)
            print(f"wrote {OUTPUT}")
            browser.close()
    finally:
        httpd.shutdown()


if __name__ == "__main__":
    main()
