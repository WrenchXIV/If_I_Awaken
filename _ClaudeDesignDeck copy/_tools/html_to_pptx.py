"""
Convert site-deck v6.html into an editable PowerPoint (.pptx).

Each slide section is rendered in a real browser at native 1920x1080,
then every visible text node, image, and shaped div is captured with its
on-screen bounding box and emitted as a separate PPTX shape so it can be
moved/resized/edited in Google Slides or PowerPoint.

Run:
    python3 html_to_pptx.py [--slide INDEX] [--all]
"""
from __future__ import annotations

import argparse
import http.server
import io
import os
import socketserver
import sys
import threading
import time
import urllib.request
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ----------------------- Constants -----------------------
DECK_DIR = Path(__file__).resolve().parent.parent  # the "_ClaudeDesignDeck copy" folder
DECK_FILE = "site-deck v6.html"
OUTPUT = DECK_DIR / "exports" / "site-deck-v6.pptx"

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
EMU_PER_PX = 12192000 / 1920  # = 6350 EMU per pixel (PPTX is 13.333" x 7.5")

# Map web fonts to their nearest PowerPoint/Google-Slides available equivalents.
FONT_MAP = {
    "Oswald": "Impact",
    "Impact": "Impact",
    "Haettenschweiler": "Impact",
    "Inter": "Arial",
    "Inter Tight": "Arial",
    "system-ui": "Arial",
    "Playfair Display": "Playfair Display",  # available in Google Slides
    "Barlow Condensed": "Oswald",
    "Fraunces": "Playfair Display",
    "Cormorant Garamond": "Cormorant Garamond",
    "EB Garamond": "EB Garamond",
    "Georgia": "Georgia",
}


def map_font(family_str: str) -> str:
    """Pick the first font in a CSS font-family list we have a mapping for."""
    if not family_str:
        return "Arial"
    for raw in family_str.split(","):
        name = raw.strip().strip("'").strip('"')
        if name in FONT_MAP:
            return FONT_MAP[name]
    # fallback: first listed font, sans-serif default
    first = family_str.split(",")[0].strip().strip("'").strip('"')
    return first or "Arial"


def parse_rgb(rgb_str: str) -> RGBColor | None:
    """Parse a CSS rgb()/rgba() string into RGBColor; None if transparent or invalid."""
    if not rgb_str:
        return None
    s = rgb_str.strip()
    if s.startswith("rgba"):
        nums = s[s.index("(") + 1 : s.rindex(")")].split(",")
        if len(nums) == 4:
            try:
                a = float(nums[3])
            except ValueError:
                a = 1.0
            if a < 0.05:
                return None
            r, g, b = (int(float(n.strip())) for n in nums[:3])
            return RGBColor(r, g, b)
    elif s.startswith("rgb"):
        nums = s[s.index("(") + 1 : s.rindex(")")].split(",")
        if len(nums) >= 3:
            r, g, b = (int(float(n.strip())) for n in nums[:3])
            return RGBColor(r, g, b)
    return None


def px_to_emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


# ----------------------- Local server -----------------------
class _QuietHandler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, format, *args):  # noqa
        pass


def start_server(directory: Path, port: int = 8765) -> socketserver.TCPServer:
    os.chdir(directory)
    socketserver.TCPServer.allow_reuse_address = True
    httpd = socketserver.TCPServer(("127.0.0.1", port), _QuietHandler)
    t = threading.Thread(target=httpd.serve_forever, daemon=True)
    t.start()
    return httpd


# ----------------------- Extractor (runs in browser) -----------------------
EXTRACTOR_JS = r"""
() => {
  // Walks the currently-visible slide element and returns a flat list
  // of editable items: text (one entry per visible text run, located via
  // Range so display:contents wrappers don't break us), images,
  // background-image divs, and colored shapes.

  const stage = document.querySelector('deck-stage');
  if (!stage) return { error: 'no deck-stage' };

  const slides = stage.querySelectorAll('section.slide');
  let active = null;
  for (const s of slides) {
    const cs = getComputedStyle(s);
    if (cs.display === 'none' || cs.visibility === 'hidden') continue;
    const r = s.getBoundingClientRect();
    if (r.width > 0 && r.height > 0) { active = s; break; }
  }
  if (!active) return { error: 'no active slide' };

  const stageRect = stage.getBoundingClientRect();
  const innerW = parseFloat(stage.getAttribute('width')) || 1920;
  const scale = stageRect.width / innerW;

  const items = [];

  function rectToBox(r) {
    return {
      x: (r.left - stageRect.left) / scale,
      y: (r.top - stageRect.top) / scale,
      w: r.width / scale,
      h: r.height / scale,
    };
  }

  function isHiddenChain(el) {
    let p = el;
    while (p && p !== document.body) {
      const cs = getComputedStyle(p);
      if (cs.display === 'none' || cs.visibility === 'hidden' || parseFloat(cs.opacity) === 0) {
        return true;
      }
      p = p.parentElement;
    }
    return false;
  }

  // Slide background
  const slideCS = getComputedStyle(active);
  items.push({
    kind: 'slide-bg',
    bg: slideCS.backgroundColor,
    bgImage: slideCS.backgroundImage,
    color: slideCS.color,
  });

  // ---- IMAGES ----
  const imageEls = new Set();
  active.querySelectorAll('img').forEach(img => {
    if (isHiddenChain(img)) return;
    const r = img.getBoundingClientRect();
    if (r.width < 2 || r.height < 2) return;
    const b = rectToBox(r);
    items.push({
      kind: 'image',
      src: img.currentSrc || img.src,
      alt: img.alt || '',
      ...b,
    });
    imageEls.add(img);
  });

  // ---- BACKGROUND-IMAGE divs ----
  const bgImgEls = new Set();
  active.querySelectorAll('*').forEach(el => {
    if (isHiddenChain(el) || el === active) return;
    const cs = getComputedStyle(el);
    const bgi = cs.backgroundImage;
    if (!bgi || bgi === 'none' || !bgi.includes('url(')) return;
    const r = el.getBoundingClientRect();
    if (r.width < 20 || r.height < 20) return;
    const url = (bgi.match(/url\(["']?([^"')]+)["']?\)/) || [])[1] || '';
    items.push({
      kind: 'bg-image',
      src: url,
      bg: cs.backgroundColor,
      ...rectToBox(r),
    });
    bgImgEls.add(el);
  });

  // ---- TEXT (one entry per logical text run, via TreeWalker on text nodes) ----
  // Group adjacent text nodes that share the same nearest styled element.
  // This avoids 1-entry-per-word and gives us natural-sized text boxes.
  function nearestStyleAnchor(textNode) {
    let p = textNode.parentElement;
    while (p && p !== active) {
      // Anchor on elements that introduce their own typography or layout.
      const cs = getComputedStyle(p);
      if (cs.display !== 'inline' && cs.display !== 'contents') return p;
      // anchor on elements with their own font-size override
      p = p.parentElement;
    }
    return textNode.parentElement || active;
  }

  const blockGroups = new Map(); // anchor element -> { texts:[], rect, anchor }

  const walker = document.createTreeWalker(active, NodeFilter.SHOW_TEXT, {
    acceptNode(n) {
      if (!n.nodeValue || !n.nodeValue.trim()) return NodeFilter.FILTER_REJECT;
      const tag = n.parentElement && n.parentElement.tagName.toLowerCase();
      if (tag === 'script' || tag === 'style') return NodeFilter.FILTER_REJECT;
      return NodeFilter.FILTER_ACCEPT;
    },
  });

  let tn;
  while ((tn = walker.nextNode())) {
    if (isHiddenChain(tn.parentElement)) continue;
    // Range-based bbox works regardless of display:contents/inline.
    const range = document.createRange();
    range.selectNodeContents(tn);
    const r = range.getBoundingClientRect();
    if (r.width < 2 || r.height < 2) continue;

    const anchor = nearestStyleAnchor(tn);
    const txt = tn.nodeValue.replace(/\s+/g, ' ').trim();

    let g = blockGroups.get(anchor);
    if (!g) {
      g = {
        anchor,
        texts: [],
        left: r.left, top: r.top, right: r.right, bottom: r.bottom,
      };
      blockGroups.set(anchor, g);
    } else {
      g.left = Math.min(g.left, r.left);
      g.top = Math.min(g.top, r.top);
      g.right = Math.max(g.right, r.right);
      g.bottom = Math.max(g.bottom, r.bottom);
    }
    g.texts.push(txt);
  }

  for (const g of blockGroups.values()) {
    const cs = getComputedStyle(g.anchor);
    const text = g.texts.join(' ').replace(/\s+/g, ' ').trim();
    if (!text) continue;
    const b = rectToBox({
      left: g.left, top: g.top,
      right: g.right, bottom: g.bottom,
      width: g.right - g.left, height: g.bottom - g.top,
    });
    items.push({
      kind: 'text',
      text,
      ...b,
      fontFamily: cs.fontFamily,
      fontSize: parseFloat(cs.fontSize),
      fontWeight: cs.fontWeight,
      fontStyle: cs.fontStyle,
      color: cs.color,
      textAlign: cs.textAlign,
      letterSpacing: cs.letterSpacing,
      textTransform: cs.textTransform,
      lineHeight: cs.lineHeight,
    });
  }

  // ---- SHAPES (visible background-color rectangles that aren't already images/bg-images) ----
  active.querySelectorAll('*').forEach(el => {
    if (isHiddenChain(el)) return;
    if (el === active) return;
    if (imageEls.has(el) || bgImgEls.has(el)) return;
    const cs = getComputedStyle(el);
    const bg = cs.backgroundColor;
    const m = bg.match(/rgba?\(([^)]+)\)/);
    if (!m) return;
    const parts = m[1].split(',').map(s => parseFloat(s));
    const a = parts.length >= 4 ? parts[3] : 1;
    if (a < 0.1) return;
    const r = el.getBoundingClientRect();
    if (r.width < 6 || r.height < 6) return;
    // skip slide-sized containers
    if (r.width >= stageRect.width * 0.98 && r.height >= stageRect.height * 0.98) return;
    items.push({
      kind: 'shape',
      bg,
      ...rectToBox(r),
    });
  });

  return { items };
}
"""


# ----------------------- PPTX builder -----------------------
def add_slide(prs: Presentation, items: list[dict], slide_label: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Slide background
    bg_item = next((i for i in items if i.get("kind") == "slide-bg"), None)
    if bg_item:
        rgb = parse_rgb(bg_item.get("bg") or "")
        if rgb is not None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = rgb

    # Sort by area DESC so larger background-ish elements draw first.
    items_sorted = sorted(
        [i for i in items if i.get("kind") != "slide-bg"],
        key=lambda i: -(i.get("w", 0) * i.get("h", 0)),
    )

    for it in items_sorted:
        try:
            kind = it["kind"]
            x = px_to_emu(it["x"])
            y = px_to_emu(it["y"])
            w = px_to_emu(max(it["w"], 4))
            h = px_to_emu(max(it["h"], 4))

            # Clip to slide area
            if x < 0:
                w += x
                x = 0
            if y < 0:
                h += y
                y = 0
            if w <= 0 or h <= 0:
                continue

            if kind == "image":
                _add_image(slide, it["src"], x, y, w, h)
            elif kind == "bg-image":
                _add_image(slide, it["src"], x, y, w, h)
            elif kind == "shape":
                _add_shape(slide, it, x, y, w, h)
            elif kind == "text":
                _add_text(slide, it, x, y, w, h)
        except Exception as e:
            print(f"  ! skipped {it.get('kind')}: {e}", file=sys.stderr)


def _add_image(slide, src: str, x: int, y: int, w: int, h: int) -> None:
    if not src:
        return
    if src.startswith("data:"):
        return  # skip inline data-uris for now
    # Resolve relative to local server / file
    if src.startswith("http://127.0.0.1") or src.startswith("http://localhost"):
        data = urllib.request.urlopen(src, timeout=10).read()
    elif src.startswith(("http://", "https://")):
        return  # skip external
    else:
        p = (DECK_DIR / src.lstrip("/")).resolve()
        if not p.exists():
            return
        data = p.read_bytes()
    # python-pptx can't embed WEBP — transcode to PNG when needed.
    lower = src.lower().split("?")[0]
    if lower.endswith(".webp") or data[:4] == b"RIFF":
        try:
            from PIL import Image
            im = Image.open(io.BytesIO(data)).convert("RGBA")
            buf = io.BytesIO()
            im.save(buf, format="PNG")
            data = buf.getvalue()
        except Exception:
            return
    slide.shapes.add_picture(io.BytesIO(data), x, y, width=w, height=h)


def _add_shape(slide, it: dict, x: int, y: int, w: int, h: int) -> None:
    rgb = parse_rgb(it.get("bg") or "")
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()  # no border
    if rgb is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb
    else:
        shp.fill.background()


def _add_text(slide, it: dict, x: int, y: int, w: int, h: int) -> None:
    # Add a small horizontal margin so text doesn't get cropped right on the edge.
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    text = it["text"]
    p = tf.paragraphs[0]
    align = it.get("textAlign", "left")
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align in ("justify", "justify-all"):
        p.alignment = PP_ALIGN.JUSTIFY

    if it.get("textTransform") == "uppercase":
        text = text.upper()

    run = p.add_run()
    run.text = text

    font = run.font
    font.name = map_font(it.get("fontFamily", ""))
    fs_px = float(it.get("fontSize") or 16)
    # Browser CSS px → PowerPoint pt. The deck is designed at 1920x1080;
    # most slide text sizes look right at face value as points (rough but ok).
    font.size = Pt(max(8, fs_px * 0.75))
    weight = it.get("fontWeight") or "400"
    try:
        font.bold = int(weight) >= 600
    except ValueError:
        font.bold = weight in ("bold", "bolder")
    font.italic = it.get("fontStyle") == "italic"
    rgb = parse_rgb(it.get("color") or "")
    if rgb is not None:
        font.color.rgb = rgb


# ----------------------- Main -----------------------
def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--slide", type=int, default=None,
                    help="Convert just one slide (0-indexed)")
    ap.add_argument("--all", action="store_true",
                    help="Convert all slides")
    args = ap.parse_args()
    if args.slide is None and not args.all:
        args.slide = 0  # default: just the cover

    httpd = start_server(DECK_DIR)
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            ctx = browser.new_context(viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX})
            page = ctx.new_page()
            url = f"http://127.0.0.1:8765/{DECK_FILE.replace(' ', '%20')}"
            page.goto(url, wait_until="networkidle")
            page.wait_for_function("document.querySelector('deck-stage') && document.querySelector('deck-stage')._slides")

            total = page.evaluate("document.querySelector('deck-stage')._slides.length")
            print(f"deck-stage reports {total} slides")

            if args.all:
                indices = list(range(total))
            else:
                indices = [args.slide]

            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)

            for i in indices:
                page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
                page.wait_for_timeout(250)  # let it settle
                page.wait_for_load_state("networkidle")
                result = page.evaluate(EXTRACTOR_JS)
                if "error" in result:
                    print(f"slide {i}: {result['error']}")
                    continue
                items = result["items"]
                print(f"slide {i:02d}: extracted {len(items)} items")
                if os.environ.get("DUMP"):
                    import json
                    print(json.dumps(items, indent=2)[:4000])
                add_slide(prs, items, slide_label=f"Slide {i+1}")

            OUTPUT.parent.mkdir(parents=True, exist_ok=True)
            prs.save(OUTPUT)
            print(f"wrote {OUTPUT}")
            browser.close()
    finally:
        httpd.shutdown()


if __name__ == "__main__":
    main()
